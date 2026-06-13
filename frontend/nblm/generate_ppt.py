import os
import sys

def install_and_generate():
    # Attempt to import python-pptx. If not found, attempt to install it.
    try:
        import pptx
    except ImportError:
        print("python-pptx is not installed. Installing it via pip...")
        import subprocess
        try:
            subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"], check=True)
            import pptx
            print("python-pptx successfully installed!")
        except Exception as e:
            print(f"Error installing python-pptx: {e}")
            print("Please install python-pptx manually by running 'pip install python-pptx'.")
            return
            
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    # 1. Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    # Color definitions (Whiteboard palette)
    c_white = RGBColor(255, 255, 255)
    c_dark = RGBColor(30, 30, 30)
    c_blue = RGBColor(59, 130, 246)
    c_green = RGBColor(16, 185, 129)
    c_gray = RGBColor(120, 120, 120)
    c_light_gray = RGBColor(245, 245, 247)
    
    # Define slide content
    slides_data = [
        {
            "title": "50 Startups 利潤預測與分析簡報",
            "subtitle": "基於 CRISP-DM 框架與貝氏編碼技術的線性迴歸實作\n\n主講人：017 產業量化分析研究團隊\n時間：2026年6月",
            "is_cover": True
        },
        {
            "title": "1. 項目概述與商業動機 (Phase 1)",
            "bullets": [
                "商業痛點：早期新創企業財務數據稀疏，估值高度依賴主觀直覺與經驗。",
                "解決方案：藉由量化三大核心支出與設立地區，構建多元線性迴歸白盒預估模型。",
                "商業決策支援：",
                "  * 精確量化研發、行政與行銷投入的邊際利潤貢獻係數。",
                "  * 深入分析地區設立（State）對淨利潤產生的結構性效益。"
            ]
        },
        {
            "title": "2. 數據探索與相關性矩陣 (Phase 2)",
            "bullets": [
                "數據集特徵：包含 50 家具代表性的初創公司樣本（R&D, Admin, Marketing, State, Profit）。",
                "Pearson 相關係數矩陣發現：",
                "  * R&D Spend ↔ Profit (r = 0.9729)：呈現極強的正線性相關關係。",
                "  * Marketing Spend ↔ Profit (r = 0.7478)：呈現中度正線性相關。",
                "  * Administration ↔ Profit (r = 0.2007)：無顯著關聯，行政開銷並非利潤增長關鍵。"
            ]
        },
        {
            "title": "3. 類別特徵與單熱編碼 (Phase 3)",
            "bullets": [
                "「State」名目特徵包含 California, Florida, New York 三個類別。",
                "One-Hot Encoding (單熱編碼)：",
                "  * 將類別特徵展開為二元指示向量 [0, 1]。",
                "  * 虛擬變數陷阱：三個州指示向量之和等於全 1 向量，與常數項完全共線，導致設計矩陣 X^T X 不可逆。",
                "  * 解決方案：丟棄首個州別（California）作基準對照，成功排除共線性障礙。"
            ]
        },
        {
            "title": "4. 貝氏 Beta 目標編碼的推廣 (Phase 3)",
            "bullets": [
                "連續變數編碼：將 Profit 正規化至 [0, 1] 區間，設定 Beta 先驗分佈。",
                "後驗均值與方差雙特徵：同時估計與還原各州利潤後驗期望值與波動不確定性特徵。",
                "  * BTE Mean = (S_c + m * p) / (n_c + m) | BTE Var = (alpha * beta) / [(alpha + beta)^2 * (alpha + beta + 1)]",
                "本專案上 BTE 的冗餘性論證與概念驗證 (PoC) 價值：",
                "  * 冗餘度：State 僅包含 3 個州 (K=3)，使用 OHE 已極精簡，BTE 確實冗餘且過度設計。",
                "  * PoC 價值：旨在為工業界處理「高基數特徵」(如數百城市或品類 ID) 時，提供防止維度爆炸、數據洩漏與過度擬合的特徵壓縮範例。"
            ]
        },
        {
            "title": "5. 線性迴歸 OLS 求解與白盒推論 (Phase 4)",
            "bullets": [
                "最小平方法 (OLS) 解析矩陣求解公式： beta = (X^T * X)^(-1) * X^T * y",
                "統計學白盒推論：",
                "  * 估算各特徵係數的標準誤差 (Std Error) 與 t 統計量。",
                "  * 推導雙尾 p-Value，確定各投入特徵在統計學上的顯著性與信賴區間。",
                "高斯-馬可夫定理假定：模型在實作中深度檢驗參數線性、同變異與無自相關假設。"
            ]
        },
        {
            "title": "6. 模型診斷與假設檢定 (Phase 5)",
            "bullets": [
                "VIF 共線性檢測：特徵膨脹因子小於 5，無高度多重共線性，係數估計極具穩健性。",
                "殘差 Q-Q 常態檢定：殘差分位數點緊密分布於對角線，符合常態分佈，保障推論有效性。",
                "異質變異檢定：殘差 vs. 擬合值散佈圖呈均勻隨機分布，無漏斗狀趨勢，支持同變異假設。"
            ]
        },
        {
            "title": "7. 編碼模式性能對照 (Phase 5)",
            "bullets": [
                "One-Hot Encoding 迴歸：測試集 R^2 = 0.9475, MAE = $6,742.30 USD",
                "Beta Target Encoding (m=2.0) 迴歸：測試集 R^2 = 0.9482, MAE = $6,552.12 USD",
                "分析總結：BTE 模型在適度平滑下能夠有效防範噪聲，展現出更優越的泛化表現。",
                "平滑參數 m 的重要性：當 m 過大或過小時都會引發欠擬合或過度擬合，m=2.0 為本專案最佳配適。"
            ]
        },
        {
            "title": "8. 系統部署與商業結論 (Phase 6)",
            "bullets": [
                "交互式部署：基於 FastAPI + Vanilla JS 搭建手繪白板風格前後端分離預估面板。",
                "風險控制：利用 Student's t 分布與 RMSE，為用戶預測即時提供 95% 信賴預測區間。",
                "商業總結：研發支出 (R&D) 具備最高邊際利益；貝氏平滑編碼在大維度與稀疏數據下極具推廣價值；本專案完整落實白盒機器學習決策流程。"
            ]
        }
    ]
    
    # 2. Programmatically build slides
    for data in slides_data:
        slide_layout = prs.slide_layouts[6] # Blank slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide background to white
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_white
        
        # Add a sketchy visual boundary box (main card shape)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12.333)
        height = Inches(6.5)
        
        # We can add a rectangle shape for slide background card
        shape = slide.shapes.add_shape(
            pptx.enum.shapes.MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_light_gray
        shape.line.color.rgb = c_dark
        shape.line.width = Pt(2)
        
        # Add slide title
        tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.0))
        tf_title = tx_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = data["title"]
        p_title.font.name = "Arial"
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = c_dark
        
        # Add body content
        if data.get("is_cover"):
            tx_sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(4.0))
            tf_sub = tx_sub.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = data["subtitle"]
            p_sub.font.name = "Arial"
            p_sub.font.size = Pt(20)
            p_sub.font.color.rgb = c_blue
        else:
            tx_body = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
            tf_body = tx_body.text_frame
            tf_body.word_wrap = True
            
            for i, bullet in enumerate(data.get("bullets", [])):
                if i == 0:
                    p = tf_body.paragraphs[0]
                else:
                    p = tf_body.add_paragraph()
                p.text = bullet
                p.font.name = "Arial"
                p.font.size = Pt(18)
                p.font.color.rgb = c_dark
                p.space_after = Pt(12)
                
                # Check indentation
                if bullet.startswith("  *"):
                    p.level = 1
                    p.font.size = Pt(16)
                    p.font.color.rgb = c_blue
                else:
                    p.level = 0
                    
    target_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "nblm_presentation.pptx")
    prs.save(target_path)
    print(f"Presentation successfully saved to: {target_path}")

if __name__ == "__main__":
    install_and_generate()
